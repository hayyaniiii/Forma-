"""PDF conversion operations for Forma."""

import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from converters import ConversionError

PDF_MAX_WARN_BYTES = 100 * 1024 * 1024

COMPRESS_LEVELS = {
    "screen": 72,
    "ebook": 150,
    "printer": 300,
    "prepress": 300,
}


def libreoffice_available() -> bool:
    if sys.platform == "win32":
        paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        if any(os.path.isfile(p) for p in paths):
            return True
    return bool(shutil.which("soffice") or shutil.which("libreoffice"))


def weasyprint_available() -> bool:
    try:
        import weasyprint  # noqa: F401
        return True
    except ImportError:
        return False


def pdf_capabilities() -> dict:
    caps = {
        "pikepdf": _try_import("pikepdf"),
        "pdf2docx": _try_import("pdf2docx"),
        "pdf2image": _try_import("pdf2image"),
        "pptx": _try_import("pptx"),
        "docx2pdf": _try_import("docx2pdf"),
        "ebooklib": _try_import("ebooklib"),
        "pypdf": _try_import("pypdf"),
        "weasyprint": weasyprint_available(),
        "libreoffice": libreoffice_available(),
    }
    return caps


def _try_import(name: str) -> bool:
    try:
        __import__(name)
        return True
    except ImportError:
        return False


def inspect_pdf(path: Path) -> dict:
    info = {
        "encrypted": False,
        "corrupt": False,
        "image_only": False,
        "large": path.stat().st_size > PDF_MAX_WARN_BYTES,
    }
    try:
        import pikepdf
        try:
            with pikepdf.open(path) as pdf:
                _ = len(pdf.pages)
        except pikepdf.PasswordError:
            info["encrypted"] = True
            return info
    except Exception:
        info["corrupt"] = True
        return info

    if path.suffix.lower() == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            if reader.is_encrypted:
                info["encrypted"] = True
                return info
            sample = ""
            for page in reader.pages[: min(5, len(reader.pages))]:
                sample += page.extract_text() or ""
            if len(sample.strip()) < 24:
                info["image_only"] = True
        except Exception:
            pass
    return info


def _office_to_pdf(input_path: Path, output_path: Path) -> Path:
    """Word/PowerPoint → PDF via docx2pdf or LibreOffice."""
    ext = input_path.suffix.lower()
    if ext not in (".docx", ".doc", ".pptx", ".ppt"):
        raise ConversionError("Unsupported office format.")

    if _try_import("docx2pdf") and ext in (".docx", ".doc"):
        try:
            from docx2pdf import convert
            convert(str(input_path), str(output_path))
            if output_path.exists():
                return output_path
        except Exception:
            pass

    if not libreoffice_available():
        raise ConversionError(
            "LibreOffice is required for Word/PowerPoint conversion. "
            "Install LibreOffice and restart Forma."
        )

    out_dir = output_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)
    soffice = shutil.which("soffice") or r"C:\Program Files\LibreOffice\program\soffice.exe"
    if not os.path.isfile(soffice):
        for p in (
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ):
            if os.path.isfile(p):
                soffice = p
                break

    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(input_path),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
    if proc.returncode != 0:
        raise ConversionError((proc.stderr or proc.stdout or "LibreOffice failed")[:400])

    produced = out_dir / (input_path.stem + ".pdf")
    if not produced.exists():
        raise ConversionError("LibreOffice did not produce a PDF.")
    if produced != output_path:
        produced.replace(output_path)
    return output_path


def compress_pdf(input_path: Path, output_path: Path, level: str = "ebook") -> Path:
    import pikepdf
    from pikepdf import ObjectStreamMode

    try:
        with pikepdf.open(input_path) as pdf:
            pdf.save(
                output_path,
                compress_streams=True,
                object_stream_mode=ObjectStreamMode.generate,
            )
    except pikepdf.PasswordError:
        raise ConversionError("PDF is password-protected. Remove protection first.") from None
    except Exception as exc:
        raise ConversionError(f"Could not compress PDF: {exc}") from exc

    _ = COMPRESS_LEVELS.get(level, 150)
    return output_path


def estimate_compress_size(input_path: Path, level: str) -> float:
    ratio = {"screen": 0.35, "ebook": 0.55, "printer": 0.75, "prepress": 0.9}.get(
        level, 0.55
    )
    return max(0.1, input_path.stat().st_size * ratio / (1024 * 1024))


def pdf_to_docx(input_path: Path, output_path: Path, language: str = "en") -> Path:
    if not _try_import("pdf2docx"):
        raise ConversionError("pdf2docx is not installed.")
    try:
        from pdf2docx import Converter
        cv = Converter(str(input_path))
        cv.convert(str(output_path))
        cv.close()
    except Exception as exc:
        msg = str(exc).lower()
        if "password" in msg or "encrypted" in msg:
            raise ConversionError("PDF is password-protected.") from exc
        raise ConversionError(f"PDF to Word failed: {exc}") from exc
    if not output_path.exists():
        raise ConversionError("Conversion produced no output.")
    return output_path


def pdf_to_pptx(input_path: Path, output_path: Path, slide_per_page: bool = True) -> Path:
    if not _try_import("pdf2image") or not _try_import("pptx"):
        raise ConversionError("pdf2image and python-pptx are required.")
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches

    try:
        images = convert_from_path(str(input_path), dpi=150)
    except Exception as exc:
        raise ConversionError(
            f"Could not rasterize PDF (is Poppler installed?): {exc}"
        ) from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for img in images if slide_per_page else images[:1]:
        slide = prs.slides.add_slide(blank)
        tmp = Path(tempfile.mkdtemp()) / "slide.png"
        img.save(tmp, "PNG")
        slide.shapes.add_picture(str(tmp), 0, 0, width=prs.slide_width)
        tmp.unlink(missing_ok=True)

    prs.save(output_path)
    return output_path


def epub_to_pdf(
    input_path: Path,
    output_path: Path,
    page_size: str = "A4",
    font_size: int = 14,
    margins: str = "normal",
) -> Path:
    if not _try_import("ebooklib"):
        raise ConversionError("ebooklib is not installed.")

    import ebooklib
    from ebooklib import epub

    margin_css = {
        "normal": "2cm",
        "narrow": "1cm",
        "wide": "3cm",
    }.get(margins, "2cm")

    size_css = {
        "A4": "A4",
        "Letter": "letter",
        "A5": "A5",
    }.get(page_size, "A4")

    book = epub.read_epub(str(input_path))
    html_parts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        html_parts.append(item.get_content().decode("utf-8", errors="ignore"))

    combined = f"""<!DOCTYPE html><html><head>
    <style>
    @page {{ size: {size_css}; margin: {margin_css}; }}
    body {{ font-family: Georgia, serif; font-size: {font_size}px; line-height: 1.5; }}
    </style></head><body>{"".join(html_parts)}</body></html>"""

    if weasyprint_available():
        from weasyprint import HTML
        HTML(string=combined).write_pdf(str(output_path))
        return output_path

    # Fallback: simple text extraction to PDF via reportlab if available
    try:
        from reportlab.lib.pagesizes import A4, letter
        from reportlab.pdfgen import canvas
        from bs4 import BeautifulSoup

        sizes = {"A4": A4, "Letter": letter, "A5": A4}
        c = canvas.Canvas(str(output_path), pagesize=sizes.get(page_size, A4))
        text = BeautifulSoup(combined, "html.parser").get_text("\n")
        y = 750
        for line in text.splitlines()[:2000]:
            if y < 50:
                c.showPage()
                y = 750
            c.drawString(50, y, line[:120])
            y -= font_size + 4
        c.save()
        return output_path
    except ImportError:
        raise ConversionError(
            "EPUB to PDF requires WeasyPrint or reportlab+beautifulsoup4. "
            "Run: py -m pip install weasyprint"
        ) from None


def pdf_to_epub(input_path: Path, output_path: Path) -> Path:
    if not _try_import("pypdf") or not _try_import("ebooklib"):
        raise ConversionError("pypdf and ebooklib are required.")
    from pypdf import PdfReader
    import ebooklib
    from ebooklib import epub

    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        raise ConversionError("PDF is password-protected.")

    book = epub.EpubBook()
    book.set_identifier("forma-convert")
    book.set_title(input_path.stem)
    book.set_language("en")
    chapters = []
    spine = ["nav"]

    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip() or f"[Page {i + 1} — no extractable text]"
        ch = epub.EpubHtml(
            title=f"Page {i + 1}",
            file_name=f"page_{i + 1}.xhtml",
            lang="en",
        )
        ch.content = f"<h1>Page {i + 1}</h1><pre>{text}</pre>"
        book.add_item(ch)
        chapters.append(ch)
        spine.append(ch)

    book.toc = chapters
    book.spine = spine
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(output_path), book)
    return output_path


def run_pdf_operation(
    operation: str,
    input_path: Path,
    output_dir: Path,
    options: dict | None = None,
) -> Path:
    options = options or {}
    stem = input_path.stem
    output_dir.mkdir(parents=True, exist_ok=True)

    info = inspect_pdf(input_path) if input_path.suffix.lower() == ".pdf" else {}
    if info.get("encrypted"):
        raise ConversionError("PDF is password-protected. Remove protection first.")
    if info.get("corrupt"):
        raise ConversionError("File appears corrupt or is not a valid PDF.")

    if operation == "compress":
        out = output_dir / f"{stem}_compressed.pdf"
        return compress_pdf(input_path, out, options.get("level", "ebook"))

    if operation == "to-docx":
        out = output_dir / f"{stem}.docx"
        return pdf_to_docx(input_path, out, options.get("language", "en"))

    if operation == "to-pptx":
        out = output_dir / f"{stem}.pptx"
        return pdf_to_pptx(
            input_path, out, options.get("slide_per_page", True)
        )

    if operation == "from-docx":
        out = output_dir / f"{stem}.pdf"
        return _office_to_pdf(input_path, out)

    if operation == "from-pptx":
        out = output_dir / f"{stem}.pdf"
        return _office_to_pdf(input_path, out)

    if operation == "epub-to-pdf":
        out = output_dir / f"{stem}.pdf"
        return epub_to_pdf(
            input_path,
            out,
            page_size=options.get("page_size", "A4"),
            font_size=int(options.get("font_size", 14)),
            margins=options.get("margins", "normal"),
        )

    if operation == "pdf-to-epub":
        out = output_dir / f"{stem}.epub"
        return pdf_to_epub(input_path, out)

    raise ConversionError(f"Unknown PDF operation: {operation}")
